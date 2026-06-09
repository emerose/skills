"""Shared, deterministic readers for CRO source files → tidy rows.

Each reader returns (header: list[str], rows: list[list]) so a recipe can write a
CSV and the extractor can record provenance. Readers are pure functions of the
source bytes — re-running on the same input yields identical output (diffable).

Formats: .xlsx (openpyxl) and GraphPad Prism cover the large majority of measurement
data. Prism comes in three on-disk shapes; the Prism readers sniff and route by the
file's leading bytes (see `read_pzfx` / `_prism_sniff`):
  - XML  `.pzfx`  (starts `<?xml`)         — parsed here directly;
  - zip  `.prism` (starts `PK\x03\x04`)    — Prism 10/11, parsed by `read_prism`;
  - binary legacy (starts `PCFFGRA4`)      — Prism 4/5/6, not parseable → clear error.
Add new readers here so recipes stay thin.
"""
from __future__ import annotations
import csv as _csv
import datetime as _dt
import io as _io
import json as _json
import xml.etree.ElementTree as ET
import zipfile as _zipfile
from pathlib import Path


def _fmt_dt(d) -> str:
    """A date/datetime cell → faithful ISO string. A datetime whose time part is
    midnight is a *date* cell (Excel stores dates as datetimes at 00:00:00) and is
    rendered date-only ('2023-05-19'); a real time component is kept
    ('2023-05-19 14:30:00'). Matches how dates are normally written by hand, so a
    faithful extraction is byte-comparable to hand-curated date columns."""
    if isinstance(d, _dt.datetime):
        if (d.hour, d.minute, d.second, d.microsecond) == (0, 0, 0, 0):
            return d.date().isoformat()
        return d.isoformat(sep=" ")
    return d.isoformat()  # datetime.date / datetime.time


# --------------------------------------------------------------------------- #
# xlsx
# --------------------------------------------------------------------------- #
def _fmt(v) -> str:
    """Stable cell → string. Integers stay integers; floats keep their value;
    date/datetime cells render as clean ISO dates; None/blank → ''. (Faithful to
    the cell value, no rounding.)"""
    if v is None:
        return ""
    if isinstance(v, (_dt.datetime, _dt.date, _dt.time)):
        return _fmt_dt(v)
    if isinstance(v, float) and v.is_integer():
        return str(int(v))
    return str(v)

def _read_xlsx_rows(path: Path, sheet: str | None) -> list[list[str]]:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb[sheet] if sheet else wb.active
    if ws is None:
        wb.close()
        raise ValueError(f"no worksheet {sheet!r} in {path}")
    rows = [[_fmt(c) for c in r] for r in ws.iter_rows(values_only=True)]
    wb.close()
    return rows


def _read_xls_rows(path: Path, sheet: str | None) -> list[list[str]]:
    """Legacy .xls via xlrd (openpyxl can't read .xls)."""
    import xlrd
    wb = xlrd.open_workbook(str(path))
    ws = wb.sheet_by_name(sheet) if sheet else wb.sheet_by_index(0)

    def cell(c) -> str:
        if c.ctype in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK):
            return ""
        if c.ctype == xlrd.XL_CELL_NUMBER:
            return str(int(c.value)) if float(c.value).is_integer() else str(c.value)
        if c.ctype == xlrd.XL_CELL_DATE:
            return _fmt_dt(xlrd.xldate.xldate_as_datetime(c.value, wb.datemode))
        return str(c.value)

    return [[cell(ws.cell(r, col)) for col in range(ws.ncols)] for r in range(ws.nrows)]


def read_xlsx_sheet(path: Path, sheet: str | None = None,
                    drop_blank_rows: bool = True) -> tuple[list[str], list[list[str]]]:
    """Dump one spreadsheet worksheet faithfully (`.xlsx` via openpyxl, `.xls` via
    xlrd): row 1 is the header, the rest are data. `drop_blank_rows` removes rows
    where every cell is blank (trailing padding / spacer rows) — it never drops a
    row that carries any value."""
    suffix = Path(path).suffix.lower()
    rows = _read_xls_rows(path, sheet) if suffix == ".xls" else _read_xlsx_rows(path, sheet)
    if not rows:
        return [], []
    header, data = rows[0], rows[1:]
    if drop_blank_rows:
        data = [r for r in data if any(c != "" for c in r)]
    return header, data

# --------------------------------------------------------------------------- #
# GraphPad Prism .pzfx
# --------------------------------------------------------------------------- #
def _tag(t: str) -> str:
    return t.split("}")[-1]

def _direct_title(col: ET.Element) -> str:
    """The column's own <Title> (direct child), not a nested one."""
    for ch in col:
        if _tag(ch.tag) == "Title":
            return "".join(ch.itertext()).strip()
    return ""

def _subcolumns(col: ET.Element) -> list[list[str]]:
    out = []
    for sc in col:
        if _tag(sc.tag) == "Subcolumn":
            out.append([(d.text or "").strip() for d in sc if _tag(d.tag) == "d"])
    return out

def read_pzfx_structured(path: Path):
    """Structured view of each table: (table_title, x_values, [(y_title, [subcolumns])]).
    `x_values` is the first X-family subcolumn; each Y column keeps its subcolumns
    (replicates) intact — the basis for a tidy long-format reshape.

    Sniffs the file: a zip `.prism` routes to `read_prism_structured`, the legacy
    binary format raises, and XML `.pzfx` is parsed here."""
    fmt = _prism_sniff(path)
    if fmt == "zip":
        return read_prism_structured(path)
    if fmt == "binary":
        raise ValueError(_PRISM_BINARY_MSG.format(path=Path(path).name))
    root = ET.parse(path).getroot()
    out = []
    for t in [el for el in root.iter() if _tag(el.tag) == "Table"]:
        xvals: list[str] = []
        ycols: list[tuple[str, list[list[str]]]] = []
        for col in t:
            k = _tag(col.tag)
            if k in ("XColumn", "XAdvancedColumn"):
                if not xvals:
                    subs = _subcolumns(col)
                    if subs:
                        xvals = subs[0]
            elif k == "YColumn":
                ycols.append((_direct_title(col), _subcolumns(col)))
        out.append((_direct_title(t), xvals, ycols))
    return out


def read_pzfx(path: Path) -> list[tuple[str, list[str], list[list[str]]]]:
    """Return one (table_title, header, rows) per data table, flattened wide:
    leading '' index column (0..n-1), then each X-family subcolumn as `_<k>`,
    then each Y column's subcolumns as `<col title>_<g>` where g is a running
    index across all Y subcolumns in the table. Mirrors the conventional
    pandas-style pzfx dump so the layout is recognizable and value-faithful.

    Sniffs the file: a zip `.prism` routes to `read_prism`, the legacy binary
    format raises, and XML `.pzfx` is parsed here."""
    fmt = _prism_sniff(path)
    if fmt == "zip":
        return read_prism(path)
    if fmt == "binary":
        raise ValueError(_PRISM_BINARY_MSG.format(path=Path(path).name))
    root = ET.parse(path).getroot()
    tables = [el for el in root.iter() if _tag(el.tag) == "Table"]
    results = []
    for t in tables:
        title = _direct_title(t)
        xcols, ycols = [], []
        for col in t:
            k = _tag(col.tag)
            if k in ("XColumn", "XAdvancedColumn"):
                xcols.append(col)
            elif k == "YColumn":
                ycols.append(col)
        cols: list[tuple[str, list[str]]] = []          # (name, values)
        for xc in xcols:
            for sub in _subcolumns(xc):
                cols.append((f"_{len(cols)}", sub))
        g = 0
        for yc in ycols:
            ytitle = _direct_title(yc)
            for sub in _subcolumns(yc):
                cols.append((f"{ytitle}_{g}", sub))
                g += 1
        nrows = max((len(v) for _, v in cols), default=0)
        header = [""] + [name for name, _ in cols]
        rows = []
        for i in range(nrows):
            row = [str(i)]
            for _, vals in cols:
                row.append(vals[i] if i < len(vals) else "")
            rows.append(row)
        results.append((title, header, rows))
    return results


# --------------------------------------------------------------------------- #
# GraphPad Prism .prism  (Prism 10/11 — a zip of JSON + per-table CSV)
# --------------------------------------------------------------------------- #
# A `.prism` is a zip archive (magic `PK\x03\x04`):
#   document.json                      — TOC; `sheets.data` lists input data-sheet UUIDs
#   data/sheets/<UUID>/sheet.json      — DataSheet: title + `table` (uid, dataSets, xDataSet)
#   data/sets/<UUID>.json              — DataSet: column/group title + replicate layout
#   data/tables/<table-uid>/data.csv   — the actual values, a rectangular CSV grid
# The grid's columns are, left to right: an optional leading row-titles column (present
# on grouped tables), then the X column(s), then each Y dataset's subcolumns. Cells are
# stored at full float64 precision (e.g. `72.7179749000000015`); the shortest round-trip
# float repr recovers the value Prism actually shows (`72.7179749`), so a faithful read
# is byte-identical to the same data re-exported as XML `.pzfx` (the one exception is
# survival/XY tables, where the XML export redundantly duplicates the X column — a pzfx
# quirk, not real data). Output shape matches `read_pzfx` so this is a drop-in.

_PRISM_BINARY_MSG = (
    "{path} is a legacy binary Prism file (magic 'PCFFGRA4', Prism 4/5/6 era) — its "
    "data is not stored as text and cannot be read here. Re-export from Prism as XML "
    "'.pzfx' (File ▸ Export, or save-as an older Prism format that writes XML) or as a "
    "modern '.prism', then point the recipe at that file."
)


def _prism_sniff(path: Path) -> str:
    """Classify a Prism file by its leading bytes: 'zip' (.prism), 'binary' (legacy
    PCFFGRA4), or 'xml' (.pzfx / anything else, parsed as XML)."""
    with open(path, "rb") as f:
        head = f.read(8)
    if head[:4] == b"PK\x03\x04":
        return "zip"
    if head == b"PCFFGRA4":
        return "binary"
    return "xml"


def _prism_nsub(ds: dict) -> int:
    """Number of subcolumns (replicate slots) a DataSet occupies in its table's CSV
    grid. Grouped tables store a `replicate ranges` list whose ranges (`"0~11"`, `"12"`)
    index the subcolumns, so the slot count is the max index + 1 (trailing empty slots
    included, exactly as the XML export keeps them); simpler tables carry a flat
    `replicates` list, one entry per subcolumn."""
    rr = ds.get("replicate ranges")
    if rr:
        m = 0
        for entry in rr:
            r = str(entry.get("range", "0"))
            end = int(r.split("~")[-1]) if "~" in r else int(r)
            m = max(m, end + 1)
        return m
    return len(ds.get("replicates", [])) or 1


def _prism_num(s: str) -> str:
    """A raw CSV cell → faithful string: '' for blank, the shortest round-trip repr of
    the float (integers as integers) so full-precision storage collapses back to the
    displayed value; non-numeric text is passed through verbatim."""
    s = s.strip()
    if s == "":
        return ""
    try:
        f = float(s)
    except ValueError:
        return s
    return str(int(f)) if f.is_integer() else repr(f)


def _prism_tables(path: Path):
    """Core `.prism` parser → one (table_title, [x_subcolumns], [(y_title, [y_subcolumns])])
    per input data sheet, in document order. Each subcolumn is a list of cell strings
    (one per grid row). This is the shared basis for `read_prism` (flat) and
    `read_prism_structured`, mirroring the pzfx readers' intermediate shape."""
    with _zipfile.ZipFile(path) as z:
        names = set(z.namelist())

        def js(name: str) -> dict:
            return _json.loads(z.read(name))

        doc = js("document.json")
        out = []
        for suid in doc.get("sheets", {}).get("data", []):
            sheet = js(f"data/sheets/{suid}/sheet.json")
            table = sheet.get("table") or {}
            tuid = table.get("uid")
            if not tuid:
                continue
            csv_name = f"data/tables/{tuid}/data.csv"
            if csv_name not in names:
                # Some Prism exports store a table's values in a binary `data.bin`
                # column store instead of a text `data.csv`. Only the CSV form is
                # supported; surface it rather than silently dropping the table.
                raise ValueError(
                    f"{Path(path).name}: data sheet {sheet.get('title', suid)!r} has no "
                    f"text data.csv (binary table store not supported) — re-export from Prism."
                )
            content = js(f"data/tables/{tuid}/content.json")
            ncol = int(content.get("numberOfColumns", 0))
            grid = list(_csv.reader(_io.TextIOWrapper(z.open(csv_name), encoding="utf-8")))

            def load_set(u: str) -> dict:
                return js(f"data/sets/{u}.json")

            xset = table.get("xDataSet")
            x_n = _prism_nsub(load_set(xset)) if xset else 0
            ydefs = [(ds.get("title", ""), _prism_nsub(ds))
                     for ds in (load_set(u) for u in table.get("dataSets", []))]

            # Leading row-titles columns (grouped tables) are whatever the grid carries
            # beyond the planned X + Y data columns; drop them (the pzfx readers omit
            # row titles too, generating a 0-based index instead).
            leading = max(0, ncol - (x_n + sum(n for _, n in ydefs)))

            def col(c: int) -> list[str]:
                idx = leading + c
                return [_prism_num(row[idx]) if idx < len(row) else "" for row in grid]

            x_subcols = [col(c) for c in range(x_n)]
            ycols, c = [], x_n
            for ytitle, nsub in ydefs:
                ycols.append((ytitle, [col(c + k) for k in range(nsub)]))
                c += nsub
            out.append((sheet.get("title", ""), x_subcols, ycols))
        return out


def read_prism_structured(path: Path):
    """`.prism` analogue of `read_pzfx_structured`: (table_title, x_values,
    [(y_title, [subcolumns])]) per data sheet. `x_values` is the first X subcolumn."""
    return [(title, (xs[0] if xs else []), ycols) for title, xs, ycols in _prism_tables(path)]


def read_prism(path: Path) -> list[tuple[str, list[str], list[list[str]]]]:
    """Read a modern GraphPad `.prism` (Prism 10/11, a zip of JSON + per-table CSV) into
    one (table_title, header, rows) per input data sheet — the same flattened-wide shape
    as `read_pzfx` (leading '' index column, then each X subcolumn as `_<k>`, then each Y
    column's subcolumns as `<col title>_<g>` with g running across all Y subcolumns), so
    it is a drop-in for the pzfx path. Values are faithful (full precision collapsed to
    the displayed value); empty cells are ''."""
    results = []
    for title, x_subcols, ycols in _prism_tables(path):
        cols: list[tuple[str, list[str]]] = []
        for sub in x_subcols:
            cols.append((f"_{len(cols)}", sub))
        g = 0
        for ytitle, subs in ycols:
            for sub in subs:
                cols.append((f"{ytitle}_{g}", sub))
                g += 1
        nrows = max((len(v) for _, v in cols), default=0)
        header = [""] + [name for name, _ in cols]
        rows = []
        for i in range(nrows):
            row = [str(i)]
            for _, vals in cols:
                row.append(vals[i] if i < len(vals) else "")
            rows.append(row)
        results.append((title, header, rows))
    return results


# --------------------------------------------------------------------------- #
# Word .docx (CRO study reports whose tables are the only machine-readable source)
# --------------------------------------------------------------------------- #
def read_docx_tables(path: Path) -> list[list[list[str]]]:
    """Return every table in a .docx as a list of tables, each a list of rows, each
    row a list of cell strings (text only; whitespace collapsed). Deterministic.

    Some CRO deliverables ship only as a Word report (no spreadsheet); its tables are
    then the raw source. Tables are returned in document order so a recipe can select
    by index. Merged cells repeat their text across the spanned grid positions (the
    python-docx default), which keeps every row the table's full width."""
    import docx  # provided via the engine's deps
    doc = docx.Document(str(path))
    out: list[list[list[str]]] = []
    for t in doc.tables:
        rows = []
        for row in t.rows:
            rows.append([" ".join(str(c.text).split()) for c in row.cells])
        out.append(rows)
    return out


# --------------------------------------------------------------------------- #
# PDF (CRO report tables finalized only as a PDF, e.g. Provantis exports)
# --------------------------------------------------------------------------- #
def read_pdf_pages(path: Path, pages: list[int] | None = None) -> list[list[str]]:
    """Return faithful per-page text lines from a PDF (pdfplumber, layout-aware).

    Result is a list of pages, each a list of text lines in reading order, with each
    line's internal spacing preserved (NOT collapsed) so a recipe can split a tabular
    row on whitespace. `pages` is a list of 1-based page numbers to extract (in the
    given order); None = all pages. Deterministic for given input bytes (pdfplumber's
    line grouping is a pure function of the page) — re-running yields identical output.

    Some CRO deliverables finalize their data tables only as a PDF (Provantis "Report
    Tables (for insertion)" exports etc.); that PDF is then the authoritative raw
    source. A recipe selects the relevant pages and parses the lines into rows, emitting
    with x.table(..., sources=[src]). Inspect first by printing each page's lines."""
    import pdfplumber  # provided via the engine's deps
    out: list[list[str]] = []
    with pdfplumber.open(str(path)) as pdf:
        idxs = range(len(pdf.pages)) if pages is None else [p - 1 for p in pages]
        for i in idxs:
            txt = pdf.pages[i].extract_text() or ""
            out.append(txt.splitlines())
    return out


# ---- document PROSE text (for verbatim quote-matching, e.g. claim grounding) -------
# Distinct from the *table* readers above: these return the document's flowing text,
# joined, for substring/quote checks. Deliberately NOT routed through libkit's loaders:
# its office path needs LibreOffice (soffice) on PATH and its default PDF loader uploads
# bytes to a hosted API (datalab) — a heavyweight system dep + a confidentiality risk for
# CRO deliverables — and it emits reformatted Markdown, which makes verbatim matching
# *more* flaky. These are offline, deterministic, and quote-faithful.
def read_pdf_text(path: Path) -> str:
    """All page text of a PDF, newline-joined (pdfplumber). Pure function of the bytes."""
    import pdfplumber

    with pdfplumber.open(str(path)) as pdf:
        return "\n".join((page.extract_text() or "") for page in pdf.pages)


def read_docx_text(path: Path) -> str:
    """All paragraph + table-cell text of a .docx, newline-joined (python-docx)."""
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def read_pptx_text(path: Path) -> str:
    """All deck prose, newline-joined (python-pptx). Deck text is scattered: title/body
    text frames, table cells, *grouped* shapes, and speaker notes — pull them all so a
    quote that lives in any of them is matchable."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes)
            else:
                yield shape

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in walk(slide.shapes):
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.extend(cell.text for cell in row.cells)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None:
                parts.append(notes.text)
    return "\n".join(parts)
